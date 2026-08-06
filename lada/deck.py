"""PowerPoint composition on top of the HCLTech corporate template.

Two jobs live here:

**The working template.** ``assets/ppt_template.pptx`` is a 35 MB corporate deck
carrying 202 layouts and their media. Inheriting all of that into every
generated deck is untenable, so :func:`working_template` builds a cached
derivative containing only the layouts LADA composes with - about 2 MB - by
dropping the example slides and the unreferenced layout parts.

**The builder.** :class:`DeckBuilder` adds the brand furniture the specification
asks for on top of the template's own styling: the Career Shaper wordmark on the
top-left bar, a live page-number field, the confidentiality clause, a gradient
title rule, image placeholders that Agent 3 can find and replace, presenter-note
voice-over scripts for Agent 4, and hyperlink-driven activity feedback.

Activity interactivity note
---------------------------
Feedback is implemented with slide hyperlinks onto *hidden* feedback slides
rather than hand-written ``p:timing`` trigger XML. Hyperlinks and the
``show="0"`` attribute are both narrow, well-supported constructs, whereas
malformed animation timing is the classic cause of PowerPoint's "needs repair"
prompt. Clicking an option in slideshow mode jumps to "Well done" or "Try
again", each of which links back to the activity.
"""

from __future__ import annotations

import copy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import config, graphics

# --------------------------------------------------------------------------
# Layout registry
# --------------------------------------------------------------------------
DASH = "–"  # the template uses an en dash in every layout name

LAYOUTS: dict[str, str] = {
    "cover_image": f"Cover {DASH} The Beam (Image)",
    "cover_dark": f"Cover {DASH} Progress (Dark)",
    "divider_dark": f"Divider Beam {DASH} Dark",
    "divider_light": f"Divider Beam {DASH} Light",
    "agenda": f"Agenda {DASH} Image, Table, Heavy content",
    "five_key": f"Five key points {DASH} Numbered, Image (centered)",
    "three_key": f"Three key points {DASH} Numbered, Image (centered)",
    "four_key": f"Four key points {DASH} Numbered, Image (Light)",
    "two_key": f"Two key points {DASH} Thin edge, Image right",
    "six_key": f"Six key points {DASH} Numbered, Image (Light)",
    "recap": f"Five key points {DASH} Gradient Top (Light)",
    "title_content": "Title and Content",
    "title_content_sub": "Title and Content w/t Subtitle",
    "title_only": "Title Only",
    "blank": "Blank",
    "table": "Table",
}

#: Layouts whose background is a dark solid - these need reversed branding.
DARK_LAYOUT_KEYS = {"cover_dark", "divider_dark"}

TEMPLATE_SOURCE = config.ASSETS_DIR / "ppt_template.pptx"
WORKING_TEMPLATE = config.ASSETS_DIR / "lada_working_template.pptx"

#: Shape-name prefix Agent 3 looks for when swapping in generated artwork.
IMAGE_SLOT_PREFIX = "LADA_IMAGE"
#: Shape-name prefix Agent 4 uses for the embedded narration badge.
AUDIO_SLOT_PREFIX = "LADA_AUDIO"

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)


def working_template(force: bool = False) -> Path:
    """Build (and cache) the slimmed template LADA composes with."""
    if WORKING_TEMPLATE.exists() and not force:
        if (not TEMPLATE_SOURCE.exists()
                or WORKING_TEMPLATE.stat().st_mtime >= TEMPLATE_SOURCE.stat().st_mtime):
            return WORKING_TEMPLATE
    if not TEMPLATE_SOURCE.exists():
        raise FileNotFoundError(
            f"Corporate template not found at {TEMPLATE_SOURCE}. Place "
            "'ppt_template.pptx' in the assets folder."
        )

    prs = Presentation(str(TEMPLATE_SOURCE))

    # Drop the sample slides; their media is what makes the file huge.
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.get(qn("r:id")))
        slide_ids.remove(slide_id)

    keep = set(LAYOUTS.values())
    for master in prs.slide_masters:
        layout_ids = master.element.get_or_add_sldLayoutIdLst()
        for layout_id in list(layout_ids):
            rel_id = layout_id.get(qn("r:id"))
            try:
                layout = master.part.rels[rel_id].target_part.slide_layout
            except (KeyError, AttributeError):
                continue
            if layout.name not in keep:
                master.part.drop_rel(rel_id)
                layout_ids.remove(layout_id)

        # The master already carries a top-right "Career Shaper(TM)" wordmark.
        # LADA places the actual logo on the top-left bar as specified, so drop
        # the text version rather than ship the brand twice on every slide.
        for shape in list(master.shapes):
            if shape.is_placeholder or not shape.has_text_frame:
                continue
            if "career shaper" in shape.text_frame.text.strip().lower():
                _delete_shape(shape)

    WORKING_TEMPLATE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(WORKING_TEMPLATE))
    return WORKING_TEMPLATE


# --------------------------------------------------------------------------
# Small XML helpers
# --------------------------------------------------------------------------
def _add_slide_number_field(text_frame, colour: str, size: int = 9) -> None:
    """Insert a live ``<a:fld type="slidenum">`` so numbering stays correct."""
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    run_xml = (
        '<a:fld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'id="{B7A8C1D2-0000-4000-8000-0000000000AA}" type="slidenum">'
        f'<a:rPr lang="en-IN" sz="{size * 100}" b="1">'
        f'<a:solidFill><a:srgbClr val="{colour}"/></a:solidFill>'
        '<a:latin typeface="Segoe UI"/></a:rPr>'
        "<a:t>1</a:t></a:fld>"
    )
    from lxml import etree
    paragraph._p.append(etree.fromstring(run_xml))


#: Notes-page geometry used when the placeholder has to be created from scratch.
_NOTES_BODY_XML = """
<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="2" name="Notes Placeholder 1"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" idx="1"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="685800" y="4343400"/>
      <a:ext cx="5486400" cy="4114800"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr vert="horz" wrap="square" lIns="91440" tIns="45720" rIns="91440"
              bIns="45720" anchor="t"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="en-IN" sz="1200"/></a:p>
  </p:txBody>
</p:sp>
"""


def ensure_notes_text_frame(slide):
    """Return the notes text frame, creating the placeholder if it is missing.

    The corporate template ships an **empty** notes master, so python-pptx finds
    no body placeholder to clone and ``notes_text_frame`` is ``None``. Presenter
    notes are not optional here - they carry the voice-over script Agent 4
    narrates - so build the placeholder with explicit geometry when absent.
    """
    notes_slide = slide.notes_slide
    frame = notes_slide.notes_text_frame
    if frame is not None:
        return frame

    from lxml import etree
    sp_tree = notes_slide.shapes._spTree
    sp_tree.append(etree.fromstring(_NOTES_BODY_XML.strip()))

    frame = notes_slide.notes_text_frame
    if frame is None:  # pragma: no cover - defensive
        raise RuntimeError("Could not create a presenter-notes placeholder.")
    return frame


def set_hidden(slide) -> None:
    """Mark a slide hidden: skipped on advance, still reachable by hyperlink."""
    slide._element.set("show", "0")


def is_hidden(slide) -> bool:
    return slide._element.get("show") == "0"


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


# --------------------------------------------------------------------------
# Builder
# --------------------------------------------------------------------------
@dataclass
class ImageSlot:
    """A picture placeholder awaiting Agent 3."""
    slide_index: int
    slide_number: int
    shape_name: str
    left: int
    top: int
    width: int
    height: int
    prompt: str
    concept: str
    chips: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "slide_index": self.slide_index,
            "slide_number": self.slide_number,
            "shape_name": self.shape_name,
            "left": self.left, "top": self.top,
            "width": self.width, "height": self.height,
            "prompt": self.prompt, "concept": self.concept,
            "chips": self.chips,
        }


class DeckBuilder:
    """Composes one branded deck."""

    def __init__(self, palette: dict[str, str], *, deck_title: str,
                 entity: str = "", confidentiality: str | None = None):
        self.palette = palette
        self.prs = Presentation(str(working_template()))
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.deck_title = deck_title
        self.entity = entity
        self.confidentiality = confidentiality or config.CONFIDENTIALITY_NOTE
        self.image_slots: list[ImageSlot] = []
        self._layouts = {layout.name: layout for layout in self.prs.slide_layouts}
        self._logo = graphics.logo_compact(palette)
        self._logo_dark = graphics.logo_compact(palette, reversed_=True)
        self._placeholder_cache: dict[str, Path] = {}

    def _placeholder_png(self, width: int | None = None,
                         height: int | None = None) -> Path:
        """Dashed placeholder frame matched to the slot's aspect ratio.

        ``insert_picture`` centre-crops to fill, so a single landscape frame gets
        its caption sliced off in the tall picture columns several layouts use.
        """
        ratio = 16 / 9
        if width and height and height > 0:
            ratio = max(0.25, min(4.0, width / height))
        key = f"{ratio:.2f}"
        cached = self._placeholder_cache.get(key)
        if cached and cached.exists():
            return cached
        base = 900
        if ratio >= 1:
            size = (base, max(240, int(base / ratio)))
        else:
            size = (max(240, int(base * ratio)), base)
        path = config.ASSETS_DIR / (
            f"image_placeholder_{self.palette['primary_blue'].lstrip('#')}"
            f"_{key.replace('.', '-')}.png")
        if not path.exists():
            graphics.render_placeholder_frame(path, self.palette,
                                              "Concept visual", size=size)
        self._placeholder_cache[key] = path
        return path

    # ------------------------------------------------------------- palette
    def _rgb(self, key_or_hex: str) -> RGBColor:
        value = self.palette.get(key_or_hex, key_or_hex)
        return RGBColor(*config.hex_to_rgb(value))

    # --------------------------------------------------------------- slide
    def layout(self, key: str):
        name = LAYOUTS.get(key, key)
        if name not in self._layouts:
            raise KeyError(f"Layout {name!r} is not in the working template.")
        return self._layouts[name]

    def add_slide(self, layout_key: str, *, kind: str = "content"):
        """Add a slide and apply the brand furniture for its ``kind``.

        ``kind`` is one of ``cover`` (logo only - covers carry their own
        decoration and conventionally omit page numbers), ``divider``
        (logo + footer) or ``content`` (logo + title rule + footer).
        """
        slide = self.prs.slides.add_slide(self.layout(layout_key))
        self._add_furniture(slide, kind=kind,
                            dark=layout_key in DARK_LAYOUT_KEYS)
        return slide

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    @property
    def visible_count(self) -> int:
        return sum(1 for s in self.prs.slides if not is_hidden(s))

    def slide_number_of(self, slide) -> int:
        for index, candidate in enumerate(self.prs.slides, 1):
            if candidate is slide or candidate._element is slide._element:
                return index
        return 0

    # ----------------------------------------------------------- furniture
    def _add_furniture(self, slide, *, kind: str, dark: bool) -> None:
        """Career Shaper logo on the top-left bar, page number, confidentiality.

        Positions dodge the template's own furniture: the master keeps a
        "Classification: Internal" label at the very top-left and an HCLTech
        mark bottom-right, so the logo sits just below the former and the page
        number and clause run along the bottom-left.
        """
        logo = self._logo_dark if dark else self._logo
        try:
            picture = slide.shapes.add_picture(
                str(logo), Inches(0.30), Inches(0.30), height=Inches(0.34))
            picture.name = "CareerShaperLogo"
        except Exception:
            pass

        if kind == "content":
            rule = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.30), Inches(0.70),
                Inches(1.85), Pt(3.0))
            rule.name = "BrandRule"
            rule.line.fill.background()
            rule.shadow.inherit = False
            fill = rule.fill
            fill.gradient()
            fill.gradient_angle = 0
            stops = fill.gradient_stops
            stops[0].color.rgb = self._rgb("primary_purple")
            stops[0].position = 0.0
            stops[1].color.rgb = self._rgb("secondary_teal")
            stops[1].position = 1.0

        if kind == "cover":
            return  # covers carry their own decoration and no page number

        footer_colour = ("FFFFFF" if dark
                         else config.normalise_hex(
                             self.palette.get("text_muted", "#5A6B85"))[1:])

        number = slide.shapes.add_textbox(Inches(0.30), Inches(8.44),
                                          Inches(0.55), Inches(0.28))
        number.name = "SlideNumber"
        number.text_frame.word_wrap = False
        number.text_frame.margin_left = 0
        number.text_frame.margin_right = 0
        _add_slide_number_field(number.text_frame, footer_colour, size=9)

        note = slide.shapes.add_textbox(Inches(0.95), Inches(8.44),
                                        Inches(12.3), Inches(0.28))
        note.name = "ConfidentialityNote"
        frame = note.text_frame
        frame.word_wrap = False
        frame.margin_left = frame.margin_top = frame.margin_bottom = 0
        run = frame.paragraphs[0].add_run()
        run.text = f"{self.deck_title}  |  {self.confidentiality}"
        run.font.size = Pt(8)
        run.font.name = "Segoe UI"
        run.font.color.rgb = RGBColor(*config.hex_to_rgb("#" + footer_colour))

    # -------------------------------------------------------- placeholders
    @staticmethod
    def placeholder(slide, idx: int):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == idx:
                return shape
        return None

    def fill_text(self, slide, idx: int, text: str, *,
                  size: int | None = None, bold: bool | None = None,
                  colour: str | None = None,
                  bullets: Sequence[str] | None = None,
                  shrink: bool = True) -> bool:
        """Write into a layout placeholder, or drop it if there is nothing to say."""
        shape = self.placeholder(slide, idx)
        if shape is None:
            return False
        if not text and not bullets:
            _delete_shape(shape)
            return False

        frame = shape.text_frame
        frame.word_wrap = True
        lines = list(bullets) if bullets else [text]
        frame.text = str(lines[0])
        for line in lines[1:]:
            frame.add_paragraph().text = str(line)

        for paragraph in frame.paragraphs:
            if size:
                for run in paragraph.runs:
                    run.font.size = Pt(size)
            if bold is not None:
                for run in paragraph.runs:
                    run.font.bold = bold
            if colour:
                for run in paragraph.runs:
                    run.font.color.rgb = self._rgb(colour)

        if shrink:
            self._autoshrink(shape, len(" ".join(str(x) for x in lines)))
        return True

    def _autoshrink(self, shape, char_count: int) -> None:
        """Coarse font shrink so long generated text does not overflow its box.

        python-pptx cannot ask PowerPoint to autofit, so estimate capacity from
        the box area and step the size down if the text clearly will not fit.
        """
        try:
            area_in2 = (Emu(shape.width).inches * Emu(shape.height).inches)
        except Exception:
            return
        if area_in2 <= 0:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                current = run.font.size.pt if run.font.size else 18.0
                # ~ characters that fit per square inch at a given point size.
                capacity = area_in2 * (2600.0 / (current ** 2))
                if char_count > capacity > 0:
                    scaled = max(9.0, current * (capacity / char_count) ** 0.5)
                    run.font.size = Pt(round(scaled, 1))

    def drop_unused_placeholders(self, slide, keep: Iterable[int]) -> None:
        """Remove untouched placeholders so no 'Click to edit' prompts remain."""
        keep_set = set(keep)
        for shape in list(slide.placeholders):
            idx = shape.placeholder_format.idx
            if idx in keep_set:
                continue
            if shape.placeholder_format.type in (PP_PLACEHOLDER.SLIDE_NUMBER,
                                                 PP_PLACEHOLDER.FOOTER,
                                                 PP_PLACEHOLDER.DATE):
                _delete_shape(shape)
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                continue
            _delete_shape(shape)

    # -------------------------------------------------------- image slots
    def claim_image_slot(self, slide, idx: int | None, *, prompt: str,
                         concept: str, chips: Sequence[str] = (),
                         fallback_box: tuple[float, float, float, float] | None = None,
                         resize: tuple[float, float, float, float] | None = None
                         ) -> ImageSlot | None:
        """Fill a picture placeholder with the dashed frame and register it.

        Registration is what lets Agent 3 find the slot later: the shape is given
        a ``LADA_IMAGE::n`` name and the art direction is stored both in the deck
        manifest and in the shape's alt-text, so the pptx is self-describing.
        """
        slide_number = self.slide_number_of(slide)
        name = f"{IMAGE_SLOT_PREFIX}::{slide_number}"

        shape = self.placeholder(slide, idx) if idx is not None else None
        if shape is not None and resize:
            shape.left, shape.top = Inches(resize[0]), Inches(resize[1])
            shape.width, shape.height = Inches(resize[2]), Inches(resize[3])
        picture = None
        if shape is not None and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            target = (shape.left, shape.top, shape.width, shape.height)
            try:
                picture = shape.insert_picture(
                    str(self._placeholder_png(shape.width, shape.height)))
            except Exception:
                picture = None
            if picture is not None:
                # insert_picture builds a fresh p:pic with no xfrm, so it
                # re-inherits the *layout* geometry and silently discards any
                # resize applied to the slide's placeholder. Pin it explicitly.
                picture.left, picture.top, picture.width, picture.height = target
        if picture is None:
            if shape is not None:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                _delete_shape(shape)
            elif fallback_box:
                left = Inches(fallback_box[0]); top = Inches(fallback_box[1])
                width = Inches(fallback_box[2]); height = Inches(fallback_box[3])
            else:
                return None
            try:
                picture = slide.shapes.add_picture(
                    str(self._placeholder_png(width, height)),
                    left, top, width, height)
            except Exception:
                return None

        picture.name = name
        try:
            picture._element._nvXxPr.cNvPr.set("descr", prompt[:900])
        except Exception:
            pass

        slot = ImageSlot(
            slide_index=len(self.prs.slides) - 1,
            slide_number=slide_number,
            shape_name=name,
            left=int(picture.left), top=int(picture.top),
            width=int(picture.width), height=int(picture.height),
            prompt=prompt, concept=concept, chips=list(chips)[:3],
        )
        self.image_slots.append(slot)
        return slot

    # ---------------------------------------------------------- decoration
    def smart_card(self, slide, left: float, top: float, width: float,
                   height: float, *, heading: str, body: str = "",
                   number: str = "", accent: str = "primary_blue",
                   dark: bool = False, heading_size: int = 15,
                   body_size: int = 12):
        """A SmartArt-style rounded card: number chip, heading, body."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                     Inches(top), Inches(width), Inches(height))
        card.adjustments[0] = 0.06
        card.shadow.inherit = False
        card.line.color.rgb = self._rgb(accent)
        card.line.width = Pt(1.0)
        if dark:
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb(accent)
            heading_colour = body_colour = "#FFFFFF"
        else:
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            heading_colour = self.palette["dark_neutral"]
            body_colour = self.palette.get("text_muted", "#5A6B85")

        frame = card.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.20)
        frame.margin_right = Inches(0.16)
        frame.margin_top = Inches(0.13)
        frame.margin_bottom = Inches(0.10)
        frame.vertical_anchor = MSO_ANCHOR.TOP

        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        if number:
            chip = paragraph.add_run()
            chip.text = f"{number}  "
            chip.font.size = Pt(heading_size + 3)
            chip.font.bold = True
            chip.font.color.rgb = (RGBColor(0xFF, 0xFF, 0xFF) if dark
                                   else self._rgb(accent))
        head = paragraph.add_run()
        head.text = heading
        head.font.size = Pt(heading_size)
        head.font.bold = True
        head.font.name = "Segoe UI"
        head.font.color.rgb = RGBColor(*config.hex_to_rgb(heading_colour))

        if body:
            body_paragraph = frame.add_paragraph()
            body_paragraph.space_before = Pt(5)
            run = body_paragraph.add_run()
            run.text = body
            run.font.size = Pt(body_size)
            run.font.name = "Segoe UI"
            run.font.color.rgb = RGBColor(*config.hex_to_rgb(body_colour))

        self._autoshrink(card, len(heading) + len(body) + len(number))
        return card

    def chip(self, slide, left: float, top: float, width: float, height: float,
             text: str, *, fill: str = "primary_blue", text_colour: str = "#FFFFFF",
             size: int = 12, bold: bool = True, outline: bool = False,
             name: str | None = None):
        """A pill-shaped label - used for options, tags and drag tokens."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                      Inches(top), Inches(width), Inches(height))
        shape.adjustments[0] = 0.34
        shape.shadow.inherit = False
        if name:
            shape.name = name
        if outline:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.color.rgb = self._rgb(fill)
            shape.line.width = Pt(1.25)
            text_colour = self.palette["dark_neutral"]
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(fill)
            shape.line.fill.background()

        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.10)
        frame.margin_top = frame.margin_bottom = Inches(0.03)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Segoe UI"
        run.font.color.rgb = RGBColor(*config.hex_to_rgb(text_colour))
        self._autoshrink(shape, len(text))
        return shape

    def band(self, slide, top: float, height: float, text: str, *,
             fill: str = "dark_neutral", size: int = 13,
             text_colour: str = "#FFFFFF", left: float = 0.30,
             width: float = 15.4, align=PP_ALIGN.LEFT):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                      Inches(top), Inches(width), Inches(height))
        shape.adjustments[0] = 0.10
        shape.shadow.inherit = False
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill)
        shape.line.fill.background()
        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.18)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.name = "Segoe UI"
        run.font.color.rgb = RGBColor(*config.hex_to_rgb(text_colour))
        return shape

    def body_text(self, slide, left: float, top: float, width: float,
                  height: float, lines: Sequence[str], *, size: int = 13,
                  colour: str | None = None, bullet: str = "•",
                  bold_first: bool = False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        target = colour or self.palette["dark_neutral"]
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(7)
            run = paragraph.add_run()
            run.text = (f"{bullet}  {line}" if bullet else str(line))
            run.font.size = Pt(size)
            run.font.name = "Segoe UI"
            run.font.bold = bold_first and index == 0
            run.font.color.rgb = RGBColor(*config.hex_to_rgb(target))
        self._autoshrink(box, sum(len(str(x)) for x in lines))
        return box

    # ---------------------------------------------------------------- notes
    @staticmethod
    def set_notes(slide, script: str) -> None:
        """Voice-over script for Agent 4, stored in the presenter notes."""
        frame = ensure_notes_text_frame(slide)
        frame.text = (script or "").strip()

    # ----------------------------------------------------------- hyperlinks
    @staticmethod
    def link_to_slide(shape, target_slide) -> None:
        try:
            shape.click_action.target_slide = target_slide
        except Exception:
            pass

    # ------------------------------------------------------------ finalise
    def move_hidden_to_end(self) -> None:
        """Park hidden feedback slides after the visible deck.

        Interleaved hidden slides make PowerPoint's slide-number field skip
        (...7, 10, 11...) because the field counts hidden slides too. Moving
        them to the tail makes the visible run number 1..N contiguously, which
        is also what the activity answer key cites. Hyperlinks are held by
        relationship id, so reordering does not break them.
        """
        slide_ids = self.prs.slides._sldIdLst
        entries = list(slide_ids)
        flags = [is_hidden(slide) for slide in self.prs.slides]
        for element, hidden in zip(entries, flags):
            if hidden:
                slide_ids.remove(element)
                slide_ids.append(element)

    def save(self, path: Path) -> Path:
        self.move_hidden_to_end()
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        core = self.prs.core_properties
        core.title = self.deck_title
        core.author = "Career Shaper - Learning Asset Development Agent"
        core.comments = self.confidentiality
        if self.entity:
            core.category = self.entity
        self.prs.save(str(path))
        return path


# --------------------------------------------------------------------------
# Reading helpers used by Agents 3, 4 and 5
# --------------------------------------------------------------------------
def open_deck(path: Path) -> Presentation:
    return Presentation(str(path))


def image_slots(prs: Presentation) -> list[dict]:
    """Locate every registered image slot in an existing deck."""
    found: list[dict] = []
    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if (shape.name or "").startswith(IMAGE_SLOT_PREFIX):
                descr = ""
                try:
                    descr = shape._element._nvXxPr.cNvPr.get("descr") or ""
                except Exception:
                    pass
                found.append({
                    "slide_index": index,
                    "slide_number": index + 1,
                    "shape_name": shape.name,
                    "left": int(shape.left), "top": int(shape.top),
                    "width": int(shape.width), "height": int(shape.height),
                    "prompt": descr,
                })
    return found


def replace_picture(slide, shape_name: str, image_path: Path) -> bool:
    """Swap a registered placeholder picture for real artwork, keeping its box."""
    for shape in list(slide.shapes):
        if shape.name != shape_name:
            continue
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        descr = ""
        try:
            descr = shape._element._nvXxPr.cNvPr.get("descr") or ""
        except Exception:
            pass
        parent = shape._element.getparent()
        position = list(parent).index(shape._element)
        parent.remove(shape._element)
        try:
            picture = slide.shapes.add_picture(str(image_path), left, top,
                                               width, height)
        except Exception:
            return False
        picture.name = shape_name
        if descr:
            try:
                picture._element._nvXxPr.cNvPr.set("descr", descr)
            except Exception:
                pass
        # Keep z-order: pictures behind text look right, so restore the slot.
        new_element = picture._element
        new_parent = new_element.getparent()
        new_parent.remove(new_element)
        parent.insert(position, new_element)
        return True
    return False


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            parts.append(shape.text_frame.text.strip())
    return "\n".join(parts)


def slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def deck_digest(prs: Presentation, *, include_notes: bool = True,
                max_chars: int = 60_000) -> str:
    """Flatten a deck to text for the validator."""
    chunks: list[str] = []
    for index, slide in enumerate(prs.slides, 1):
        tag = " [HIDDEN FEEDBACK SLIDE]" if is_hidden(slide) else ""
        chunk = [f"--- Slide {index}{tag} ---", slide_text(slide) or "(no text)"]
        if include_notes:
            notes = slide_notes(slide)
            if notes:
                chunk.append(f"[VOICE-OVER SCRIPT] {notes}")
        chunks.append("\n".join(chunk))
    text = "\n\n".join(chunks)
    return text[:max_chars]
