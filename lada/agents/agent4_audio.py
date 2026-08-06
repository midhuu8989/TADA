"""Agent 4 - Audio Enabling Agent.

Walks every slide of every deck, reads the voice-over script Agent 2 left in the
presenter notes, and narrates it in a professional Indian-English female voice
via Gemini TTS. Each track is embedded in its own slide so the deck plays as a
self-contained e-learning asset, and a continuous full-module track is written
alongside for LMS upload.

Narration pace
--------------
Text-to-speech has no post-hoc playback-rate control: pace is a property of the
synthesis. The pace selector therefore steers the model at generation time and is
surfaced on the deck's first slide as a narration control card, so anyone opening
the deck can see which voice and pace it was built with and rerun this agent to
change it. PowerPoint's own audio control bar still gives play, pause and seek on
each slide.
"""

from __future__ import annotations

from pathlib import Path

from .. import audio as audio_embed
from .. import config, deck, llm
from ..orchestrator import AgentContext

#: Skip narrating these - hidden hyperlink feedback slides carry no script.
MIN_SCRIPT_CHARS = 8


def _fallback_script(title: str, index: int) -> str:
    if index == 1:
        return (f"Welcome to this module. We will work through {title} together, "
                "step by step.")
    return f"Next, we turn to {title}."


def _control_card(builder_palette: dict[str, str], slide, voice_label: str,
                  pace_label: str, total_seconds: float, slides: int) -> None:
    """Narration control card on slide 1, per the first-slide requirement."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30),
                                  Inches(7.55), Inches(4.60), Inches(0.80))
    card.name = "LADA_AUDIO_CONTROL"
    card.adjustments[0] = 0.14
    card.shadow.inherit = False
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(
        *config.hex_to_rgb(builder_palette["dark_neutral"]))
    card.line.color.rgb = RGBColor(
        *config.hex_to_rgb(builder_palette["secondary_teal"]))
    card.line.width = Pt(1.0)

    frame = card.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_top = Inches(0.07)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    head = frame.paragraphs[0]
    head.alignment = PP_ALIGN.LEFT
    run = head.add_run()
    run.text = "NARRATION"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.name = "Segoe UI"
    run.font.color.rgb = RGBColor(
        *config.hex_to_rgb(builder_palette["secondary_teal"]))

    minutes, seconds = divmod(int(total_seconds), 60)
    detail = frame.add_paragraph()
    detail_run = detail.add_run()
    detail_run.text = (f"{voice_label}  |  pace: {pace_label}  |  "
                       f"{slides} narrated slides  |  {minutes}m {seconds:02d}s total. "
                       "Click the speaker on any slide to play, pause or seek.")
    detail_run.font.size = Pt(9)
    detail_run.font.name = "Segoe UI"
    detail_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def run(ctx: AgentContext) -> dict:
    decks = ctx.payload(2).get("decks") or []
    if not decks:
        raise RuntimeError("Agent 2 produced no decks; rerun the deck generator.")

    voice_label = ctx.options.get("voice_label") or next(iter(config.VOICE_OPTIONS))
    voice = config.VOICE_OPTIONS.get(voice_label, config.DEFAULT_VOICE)
    pace_label = ctx.options.get("pace_label") or config.DEFAULT_PACE
    pace_style = config.NARRATION_PACE.get(pace_label,
                                           config.NARRATION_PACE[config.DEFAULT_PACE])
    accent = ctx.options.get("accent") or "warm, professional Indian English"
    autoplay = bool(ctx.options.get("autoplay", True))

    ctx.step(0.03, f"Narrating {len(decks)} deck(s) with the {voice} voice "
                   f"at {pace_label}.")

    # Pre-count narratable slides so progress is honest.
    plans: list[tuple[dict, list[tuple[int, str, str]]]] = []
    for deck_info in decks:
        path = Path(deck_info["path"])
        if not path.exists():
            continue
        presentation = deck.open_deck(path)
        titles = deck_info.get("slide_titles") or {}
        items: list[tuple[int, str, str]] = []
        for index, slide in enumerate(presentation.slides, 1):
            if deck.is_hidden(slide):
                continue
            script = deck.slide_notes(slide)
            if len(script) < MIN_SCRIPT_CHARS:
                script = _fallback_script(titles.get(str(index), "this topic"),
                                          index)
            items.append((index, script, titles.get(str(index), "")))
        plans.append((deck_info, items))

    total = sum(len(items) for _, items in plans) or 1
    done = 0
    narrated = 0
    failures: list[str] = []
    tts_blocked = False
    tts_detail = ""
    results: list[dict] = []

    for deck_index, (deck_info, items) in enumerate(plans, 1):
        path = Path(deck_info["path"])
        presentation = deck.open_deck(path)
        module_no = deck_info.get("module_sr_no", deck_index)
        audio_dir = ctx.out_dir / "audio" / f"deck{module_no}"
        audio_dir.mkdir(parents=True, exist_ok=True)

        tracks: list[dict] = []
        deck_seconds = 0.0

        for slide_number, script, title in items:
            slide = presentation.slides[slide_number - 1]
            audio_embed.strip_slide_audio(slide)
            wav_path = audio_dir / f"slide{slide_number:02d}.wav"

            if not tts_blocked:
                try:
                    ctx.client.synthesize_speech(
                        script, wav_path, voice=voice, style=pace_style,
                        accent=accent,
                        call_kind=f"tts-m{module_no}-s{slide_number}")
                except llm.QuotaError as exc:
                    tts_blocked = True
                    tts_detail = str(exc)[:280]
                    ctx.log("Text-to-speech quota exhausted for this API key; "
                            "remaining slides will keep their written scripts "
                            "without audio.", "warning")
                except llm.LLMError as exc:
                    failures.append(f"slide {slide_number}: {exc}")
                    ctx.log(f"Narration failed for slide {slide_number}: {exc}",
                            "warning")

            if wav_path.exists():
                seconds = llm.wav_duration(wav_path)
                deck_seconds += seconds
                try:
                    audio_embed.embed_slide_audio(
                        slide, wav_path, ctx.palette,
                        duration_seconds=seconds, autoplay=autoplay,
                        name=f"LADA_AUDIO::{slide_number}")
                    narrated += 1
                    tracks.append({"slide_number": slide_number, "title": title,
                                   "seconds": seconds, "words": len(script.split()),
                                   "file": wav_path.name})
                except audio_embed.AudioEmbedError as exc:
                    failures.append(f"slide {slide_number}: {exc}")

            done += 1
            if done % 2 == 0 or done == total:
                ctx.progress(0.03 + 0.90 * (done / total),
                             f"Narrated {done}/{total} slides")

        if tracks:
            _control_card(ctx.palette, presentation.slides[0], voice_label,
                          pace_label, deck_seconds, len(tracks))

        combined = None
        if tracks:
            combined = llm.concat_wavs(
                [audio_dir / t["file"] for t in tracks],
                audio_dir / f"module{module_no:02d}-full-narration.wav")
            if combined:
                ctx.register(combined, "wav",
                             f"Deck {deck_index} continuous narration")

        if tracks:
            presentation.save(str(path))
            ctx.register(path, "pptx",
                         f"Deck {deck_index} narrated: "
                         f"{deck_info.get('module_name', '')}")
        else:
            # Nothing changed - rewriting a multi-megabyte deck would only
            # churn the file and reset its timestamp.
            ctx.register(path, "pptx",
                         f"Deck {deck_index} (no narration yet): "
                         f"{deck_info.get('module_name', '')}")
        results.append({
            "filename": deck_info.get("filename"),
            "module_sr_no": module_no,
            "path": str(path),
            "narrated_slides": len(tracks),
            "total_seconds": round(deck_seconds, 1),
            "audio_dir": str(audio_dir),
            "combined_track": str(combined) if combined else "",
            "tracks": tracks,
        })
        ctx.log(f"Deck {deck_index}: {len(tracks)} slides narrated, "
                f"{deck_seconds / 60:.1f} minutes of audio.")

    total_seconds = sum(r["total_seconds"] for r in results)
    notices: list[str] = []
    if tts_blocked:
        notices.append(
            "Text-to-speech quota was exhausted part-way through this run, so "
            "some slides keep their written voice-over script but carry no audio. "
            "Rerun this agent once quota resets. "
            f"API detail: {tts_detail}")
    if failures:
        notices.append(f"{len(failures)} slide(s) could not be narrated: "
                       + "; ".join(failures[:6]))
    if autoplay:
        notices.append(
            "Narration is set to start automatically on each slide. If a viewer "
            "prefers click-to-play, select the speaker icon and set "
            "Playback > Start > On Click.")

    minutes, seconds = divmod(int(total_seconds), 60)
    covered = sum(1 for r in results if r["narrated_slides"])
    coverage = ("" if covered == len(results)
                else f" Narration covers {covered} of {len(results)} deck(s).")
    summary = (f"{narrated} slide(s) narrated across {len(results)} deck(s) - "
               f"{minutes}m {seconds:02d}s of audio in the {voice} voice "
               f"at {pace_label}.{coverage}")
    ctx.step(1.0, summary)

    return {
        "summary": summary,
        "decks": results,
        "narrated_slides": narrated,
        "total_seconds": round(total_seconds, 1),
        "voice": voice,
        "voice_label": voice_label,
        "pace_label": pace_label,
        "autoplay": autoplay,
        "tts_blocked": tts_blocked,
        "failures": failures,
        "notices": notices,
    }
