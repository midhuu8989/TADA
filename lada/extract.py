"""Text extraction from uploaded programme documentation.

Course guides arrive as PDF, Word, Excel, PowerPoint or plain text. Everything
here uses libraries LADA already depends on, so uploading a syllabus never needs
an extra install: PDFs go through PyMuPDF when present, .docx and .pptx are read
straight out of their OOXML parts, and spreadsheets go through openpyxl.
"""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

MAX_CHARS = 200_000

SUPPORTED = (".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".txt", ".md", ".csv")

_WS = re.compile(r"[ \t]+")
_BLANKS = re.compile(r"\n{3,}")


def _tidy(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _WS.sub(" ", text)
    text = _BLANKS.sub("\n\n", text)
    return text.strip()[:MAX_CHARS]


def _from_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""
    parts: list[str] = []
    with fitz.open(str(path)) as document:
        for index, page in enumerate(document, 1):
            body = page.get_text().strip()
            if body:
                parts.append(f"[page {index}]\n{body}")
    return "\n\n".join(parts)


def _from_docx(path: Path) -> str:
    """Read .docx text straight from its OOXML part.

    Paragraph ends and tabs are turned into literal whitespace markers first,
    then the run texts and those markers are collected in document order, so
    table cells and tabbed columns do not run together.
    """
    with zipfile.ZipFile(path) as archive:
        names = [n for n in archive.namelist()
                 if n == "word/document.xml"
                 or n.startswith("word/header") or n.startswith("word/footer")]
        out: list[str] = []
        for name in sorted(names):
            xml = archive.read(name).decode("utf-8", "ignore")
            xml = re.sub(r"</w:p>", "\n", xml)
            xml = re.sub(r"</w:tc>", "\t", xml)
            xml = re.sub(r"<w:(?:tab|br)[^>]*/>", "\t", xml)
            pieces = re.findall(
                r"<w:t(?:\s[^>]*)?>(.*?)</w:t>|([\n\t])", xml, re.S)
            out.append("".join(text or marker for text, marker in pieces))
    return "\n".join(out)


def _from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    parts: list[str] = []
    presentation = Presentation(str(path))
    for index, slide in enumerate(presentation.slides, 1):
        body = [shape.text_frame.text.strip() for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()]
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if body or notes:
            block = f"[slide {index}]\n" + "\n".join(body)
            if notes:
                block += f"\n[notes] {notes}"
            parts.append(block)
    return "\n\n".join(parts)


def _from_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    parts: list[str] = []
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c not in (None, "")]
                if cells:
                    rows.append(" | ".join(cells))
                if len(rows) > 400:
                    rows.append("... (truncated)")
                    break
            if rows:
                parts.append(f"[sheet: {sheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()
    return "\n\n".join(parts)


def _from_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_bytes().decode("utf-8", "ignore")


_HANDLERS = {
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".pptx": _from_pptx,
    ".xlsx": _from_xlsx,
    ".xlsm": _from_xlsx,
    ".txt": _from_text,
    ".md": _from_text,
    ".csv": _from_text,
}


def extract_text(path: Path) -> tuple[str, str]:
    """Return ``(text, note)``. ``note`` explains an empty or partial result."""
    path = Path(path)
    suffix = path.suffix.lower()
    handler = _HANDLERS.get(suffix)
    if handler is None:
        return "", (f"{suffix or 'this file type'} is not a supported document "
                    f"type. Supported: {', '.join(SUPPORTED)}.")
    try:
        text = _tidy(handler(path) or "")
    except Exception as exc:
        return "", f"Could not read {path.name}: {exc}"
    if not text:
        if suffix == ".pdf":
            return "", (f"No selectable text found in {path.name}. Scanned PDFs "
                        "need OCR before upload - paste the coverage into the "
                        "Content and Coverage boxes instead.")
        return "", f"No text could be extracted from {path.name}."
    return text, f"Extracted {len(text):,} characters from {path.name}."
